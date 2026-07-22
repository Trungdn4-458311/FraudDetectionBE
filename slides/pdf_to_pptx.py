#!/usr/bin/env python3
"""Convert a compiled PDF deck into a PPTX (one full-slide image per page).

Usage:  python slides/pdf_to_pptx.py slides/main.pdf slides/main.pptx
        (or:  pixi run slides-to-pptx)

Self-contained: PyMuPDF renders each page to a high-DPI image and python-pptx
places it on a matching 16:9 slide. This preserves the exact Beamer design; the
slides are images, not re-editable text (the editable source stays slides/main.tex).
"""
import io
import sys

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


def convert(pdf_path: str, pptx_path: str, dpi: int = 200) -> None:
    doc = fitz.open(pdf_path)
    if doc.page_count == 0:
        sys.exit(f"{pdf_path}: no pages")

    prs = Presentation()
    r0 = doc[0].rect
    aspect = r0.width / r0.height
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(round(7.5 * aspect, 3))
    blank = prs.slide_layouts[6]
    zoom = dpi / 72.0

    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
        img = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)

    prs.save(pptx_path)
    print(f"[pptx] {doc.page_count} slides -> {pptx_path} "
          f"({prs.slide_width.inches:.2f} x {prs.slide_height.inches:.2f} in @ {dpi} dpi)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        sys.exit("usage: python pdf_to_pptx.py <in.pdf> <out.pptx>")
    convert(sys.argv[1], sys.argv[2])
